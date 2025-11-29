#!/usr/bin/env python3
"""
HTML to PPTX Converter for ML-Driven Quantum Hacking Presentation
Converts presentation_20slides.html to presentation_20slides.pptx

Install requirements first:
    pip install python-pptx pillow

Usage:
    python html_to_pptx_converter.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re
import os

# Color scheme from HTML (dark theme)
COLORS = {
    'dark_bg': RGBColor(15, 23, 42),      # #0f172a
    'light_text': RGBColor(241, 245, 249), # #f1f5f9
    'secondary_text': RGBColor(203, 213, 225), # #cbd5e1
    'muted_text': RGBColor(148, 163, 184), # #94a3b8
    'accent': RGBColor(56, 189, 248),     # #38bdf8 (cyan)
    'purple': RGBColor(124, 58, 237),     # #7c3aed
    'teal': RGBColor(13, 148, 136),       # #0d9488
    'highlight': RGBColor(251, 191, 36),  # #fbbf24
    'card_bg': RGBColor(30, 41, 59),      # #1e293b
}

def extract_html_slides(html_file):
    """Extract slide content from HTML file"""
    from bs4 import BeautifulSoup
    
    with open(html_file, 'r', encoding='utf-8') as f:
        html_content = f.read()

    soup = BeautifulSoup(html_content, 'html.parser')
    slides = []

    # Find all slides using BeautifulSoup
    slide_divs = soup.find_all('div', class_='slide', attrs={'data-slide': True})

    for slide_div in slide_divs:
        slide_num = slide_div.get('data-slide')
        
        # Extract title
        title_tag = slide_div.find('h1', class_='slide-title')
        title = title_tag.get_text(strip=True) if title_tag else f"Slide {slide_num}"

        # Extract list items
        list_items = [li.get_text(strip=True) for li in slide_div.find_all('li')]
        items = [item for item in list_items if item][:12]

        # Extract equations
        equation_divs = slide_div.find_all('div', class_='equation')
        equations = [eq.get_text(strip=True) for eq in equation_divs]

        # Extract highlights
        highlight_divs = slide_div.find_all('div', class_='highlight-box')
        highlights = [h.get_text(strip=True) for h in highlight_divs if h.get_text(strip=True)]

        # Extract section titles
        section_divs = slide_div.find_all('div', class_='section-title')
        sections = [s.get_text(strip=True) for s in section_divs]

        # Extract paragraphs
        paragraphs = [p.get_text(strip=True) for p in slide_div.find_all('p') if len(p.get_text(strip=True)) > 10]

        slides.append({
            'number': int(slide_num),
            'title': title,
            'items': items,
            'equations': equations,
            'highlights': highlights,
            'sections': sections,
            'paragraphs': paragraphs
        })

    return sorted(slides, key=lambda x: x['number'])

def create_pptx_from_slides(slides_data, output_file='presentation_20slides.pptx'):
    """Create PPTX presentation from slide data"""

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print(f"Creating PPTX with {len(slides_data)} slides...")

    for slide_data in slides_data:
        # Add blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['dark_bg']

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.margin_top = Pt(10)
        title_frame.margin_left = Pt(10)
        title_frame.margin_right = Pt(10)

        p = title_frame.paragraphs[0]
        p.text = slide_data['title']
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = COLORS['light_text']

        # Add bottom border line
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(0))
        line.line.color.rgb = COLORS['purple']
        line.line.width = Pt(3)

        # Add content
        content_top = Inches(1.5)
        content_height = Inches(5.7)

        # Add slide number
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.4))
        footer_frame = footer_box.text_frame
        footer_frame.text = f"{slide_data['number']} / 20"
        footer_frame.paragraphs[0].font.size = Pt(12)
        footer_frame.paragraphs[0].font.color.rgb = COLORS['muted_text']
        footer_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Content area
        content_box = slide.shapes.add_textbox(Inches(0.7), content_top, Inches(8.6), content_height)
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        # Add equations
        for eq in slide_data['equations']:
            if text_frame.text:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            p.text = eq
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['accent']
            p.level = 0
            p.space_before = Pt(12)
            p.space_after = Pt(12)

        # Add highlights
        for highlight in slide_data['highlights']:
            if text_frame.text:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            p.text = highlight
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['secondary_text']
            p.level = 0
            p.space_before = Pt(8)
            p.space_after = Pt(8)

        # Add section titles
        for section in slide_data['sections']:
            if text_frame.text:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            p.text = section
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = COLORS['accent']
            p.level = 0
            p.space_before = Pt(10)
            p.space_after = Pt(8)

        # Add list items
        for item in slide_data['items']:
            if text_frame.text:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            p.text = item
            p.font.size = Pt(13)
            p.font.color.rgb = COLORS['secondary_text']
            p.level = 0
            p.space_before = Pt(4)
            p.space_after = Pt(4)

        # Add paragraphs
        for para in slide_data['paragraphs']:
            if text_frame.text:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            p.text = para
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['secondary_text']
            p.level = 0
            p.space_before = Pt(8)
            p.space_after = Pt(8)

        print(f"  ✓ Slide {slide_data['number']}: {slide_data['title'][:50]}")

    # Save presentation
    prs.save(output_file)
    print(f"\n✅ PPTX created successfully: {output_file}")
    print(f"   Slides: {len(slides_data)}")
    print(f"   Resolution: 10\" x 7.5\" (16:9 aspect ratio)")

    return output_file

def main():
    """Main conversion function"""
    html_file = 'presentation_20slides.html'
    output_file = 'presentation_20slides.pptx'

    if not os.path.exists(html_file):
        print(f"Error: {html_file} not found!")
        return

    print(f"Converting {html_file} to PowerPoint format...\n")

    # Extract slides
    slides_data = extract_html_slides(html_file)
    print(f"Extracted {len(slides_data)} slides\n")

    # Create PPTX
    create_pptx_from_slides(slides_data, output_file)

if __name__ == '__main__':
    main()